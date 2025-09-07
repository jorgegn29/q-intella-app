import streamlit as st
import os
from io import BytesIO
import datetime
import csv
import nltk
import docx
import PyPDF2
import yake
from wordcloud import WordCloud
import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
from sumy.parsers.plaintext import PlaintextParser
from sumy.nlp.tokenizers import Tokenizer
from sumy.summarizers.lsa import LsaSummarizer
from gtts import gTTS
from deep_translator import GoogleTranslator
from textblob import TextBlob
from fpdf import FPDF

st.set_page_config(page_title="Q-Intella - Analizador de Datos", page_icon="🟩", layout="wide")
st.markdown("""
    <style>
    .main {background-color: #e6ffe6;}
    .stApp {background-color: #e6ffe6;}
    .css-1d391kg {background-color: #e6ffe6;}
    </style>
""", unsafe_allow_html=True)

st.markdown("""
<div style='text-align:center; margin-bottom: 0.5em;'>
    <h1 style='font-size:4em; color:#228B22; font-weight:bold; text-shadow: 2px 2px 8px #b2f7b2;'>Q-Intella - Analizador de Datos</h1>
    <h2 style='font-size:2.5em; color:#228B22; font-weight:bold; text-shadow: 1px 1px 6px #b2f7b2;'>by Jorge Guillem</h2>
    <hr style='border:2px solid #228B22; width:60%; margin:auto;'>
</div>
""", unsafe_allow_html=True)
st.header("Carga de archivo")

# Inicializar historial en el estado de sesión
if "historial" not in st.session_state:
    st.session_state["historial"] = []

archivo = st.file_uploader(
    "Sube tu archivo (Word, PDF, TXT, CSV)",
    type=["docx", "pdf", "txt", "csv"],
    help="Selecciona un archivo para analizar."
)

if archivo:
    st.success("Archivo cargado correctamente.")
    st.write(f"**Nombre:** {archivo.name}")
    st.write(f"**Tamaño:** {archivo.size / 1024:.2f} KB")
    st.write(f"**Tipo:** {archivo.type}")
    # Puedes guardar el archivo si lo deseas
    # with open(os.path.join("uploads", archivo.name), "wb") as f:
    #     f.write(archivo.getbuffer())

    # --- Extracción de texto según tipo de archivo ---
    texto = ""
    if archivo.name.endswith(".txt"):
        texto = archivo.read().decode("utf-8", errors="ignore")
    elif archivo.name.endswith(".csv"):
        texto = ""
        archivo.seek(0)
        reader = csv.reader(BytesIO(archivo.read()).read().decode("utf-8", errors="ignore").splitlines())
        for row in reader:
            texto += " ".join(row) + " "
    elif archivo.name.endswith(".docx"):
        doc = docx.Document(BytesIO(archivo.read()))
        texto = " ".join([p.text for p in doc.paragraphs])
    elif archivo.name.endswith(".pdf"):
        pdf = PyPDF2.PdfReader(BytesIO(archivo.read()))
        texto = ""
        for page in pdf.pages:
            texto += page.extract_text() or ""

    # --- Extracción de palabras clave con YAKE ---
    if texto.strip():
        st.header("Palabras clave")
        kw_extractor = yake.KeywordExtractor(lan="es", n=1, top=20)
        keywords = kw_extractor.extract_keywords(texto)
        palabras = [kw[0] for kw in keywords]
        st.write("**Lista ordenada de palabras clave:**")
        st.write(palabras)

        # --- Nube de palabras ---
        st.write("**Nube de palabras:**")
        wordcloud = WordCloud(width=800, height=400, background_color="#e6ffe6", colormap="Greens").generate(" ".join(palabras))
        fig, ax = plt.subplots()
        ax.imshow(wordcloud, interpolation='bilinear')
        ax.axis('off')
        st.pyplot(fig)

    # --- Panel de estadísticas rápidas ---
    st.header("Panel de estadísticas rápidas")
    num_palabras = len(texto.split())
    num_frases = texto.count('.')
    num_caracteres = len(texto)
    st.markdown(f"**Palabras:** {num_palabras} | **Frases:** {num_frases} | **Caracteres:** {num_caracteres}")

    # --- Ranking de frases más largas/interesantes ---
    st.header("Ranking de frases más largas/interesantes")
    frases = [f.strip() for f in texto.split('.') if len(f.strip()) > 0]
    frases_ordenadas = sorted(frases, key=lambda x: len(x.split()), reverse=True)
    st.write("**Top 5 frases más largas:**")
    for i, frase in enumerate(frases_ordenadas[:5]):
        st.markdown(f"{i+1}. {frase} <span style='color:gray'>({len(frase.split())} palabras)</span>", unsafe_allow_html=True)

    # --- Análisis de sentimiento visual mejorado ---
    st.header("Análisis de sentimiento visual")
    blob = TextBlob(texto)
    sentimiento = blob.sentiment.polarity
    sentimiento_label = "NEUTRAL"
    sentimiento_color = "#FFD700"
    if sentimiento > 0.1:
        sentimiento_label = "POSITIVO"
        sentimiento_color = "#4CAF50"
    elif sentimiento < -0.1:
        sentimiento_label = "NEGATIVO"
        sentimiento_color = "#F44336"
    st.markdown(f"<div style='padding:1em; background:{sentimiento_color}; color:white; border-radius:10px; text-align:center; font-size:1.5em;'>Sentimiento: {sentimiento_label} ({sentimiento:.2f})</div>", unsafe_allow_html=True)

    # --- Frecuencia de palabras con gráfica interactiva ---
    st.header("Frecuencia de palabras (interactiva)")
    palabras_freq = pd.Series(texto.lower().split()).value_counts().head(20)
    st.bar_chart(palabras_freq)

    # --- Dashboard visual ---
    st.header("Dashboard visual de estadísticas del texto")
    palabras_freq = pd.Series(texto.lower().split()).value_counts().head(10)
    st.subheader("Frecuencia de palabras (Top 10)")
    st.bar_chart(palabras_freq)

    st.subheader("Distribución de palabras clave (queso)")
    colors = ['#4CAF50', '#81C784', '#A5D6A7', '#C8E6C9', '#388E3C']
    fig_pie, ax_pie = plt.subplots()
    ax_pie.pie([1]*len(palabras), labels=palabras, colors=colors * (len(palabras)//len(colors)+1), autopct='%1.1f%%')
    st.pyplot(fig_pie)

    st.subheader("Longitud de frases (puntos y líneas)")
    frases = [f for f in texto.split('.') if len(f.strip()) > 0]
    longitudes = [len(f.split()) for f in frases]
    fig_line, ax_line = plt.subplots()
    ax_line.plot(longitudes, marker='o', color='green', linestyle='-')
    ax_line.set_xlabel('Frase')
    ax_line.set_ylabel('Número de palabras')
    st.pyplot(fig_line)

    # --- Exportar gráficas como imágenes
    st.subheader("Exportar gráficas como imágenes")
    # Exportar nube de palabras
    wordcloud_img = BytesIO()
    fig.savefig(wordcloud_img, format='png')
    wordcloud_img.seek(0)
    st.download_button("Descargar nube de palabras (PNG)", data=wordcloud_img, file_name="nube_palabras.png", mime="image/png")
    # Exportar gráfica de pastel
    pie_img = BytesIO()
    fig_pie.savefig(pie_img, format='png')
    pie_img.seek(0)
    st.download_button("Descargar gráfica de pastel (PNG)", data=pie_img, file_name="grafica_pastel.png", mime="image/png")
    # Exportar gráfica de líneas
    line_img = BytesIO()
    fig_line.savefig(line_img, format='png')
    line_img.seek(0)
    st.download_button("Descargar gráfica de líneas (PNG)", data=line_img, file_name="grafica_lineas.png", mime="image/png")

    # --- Análisis de sentimiento ---
    st.header("Análisis de sentimiento del texto")
    blob = TextBlob(texto)
    sentimiento = blob.sentiment.polarity
    if sentimiento > 0.1:
        st.success("El texto tiene un sentimiento POSITIVO")
    elif sentimiento < -0.1:
        st.error("El texto tiene un sentimiento NEGATIVO")
    else:
        st.info("El texto tiene un sentimiento NEUTRAL")

    # --- Reconocimiento de entidades con IA ---
    st.header("Grafo interactivo de palabras clave y frases largas")
    import networkx as nx
    from pyvis.network import Network
    import tempfile
    # Crear grafo: nodos son palabras clave y frases largas
    G = nx.Graph()
    for palabra in palabras:
        G.add_node(palabra, label="Palabra clave", color="#4CAF50")
    for frase in frases_ordenadas[:5]:
        G.add_node(frase, label="Frase larga", color="#FFD700")
        # Relacionar cada frase larga con sus palabras clave presentes
        for palabra in palabras:
            if palabra in frase:
                G.add_edge(frase, palabra)
    # Visualizar grafo con pyvis
    net = Network(height="400px", width="100%", bgcolor="#e6ffe6", font_color="black")
    net.from_nx(G)
    tmp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.html')
    net.write_html(tmp_file.name)
    st.components.v1.html(open(tmp_file.name, 'r', encoding='utf-8').read(), height=400)

    st.header("Nube de palabras clave con colores aleatorios")
    import random
    color_map = ["#4CAF50", "#81C784", "#A5D6A7", "#FFD700", "#F44336", "#2196F3", "#9C27B0"]
    wordcloud = WordCloud(width=800, height=400, background_color="#e6ffe6", colormap=None,
                         color_func=lambda *args, **kwargs: random.choice(color_map)).generate(" ".join(palabras))
    fig_wc, ax_wc = plt.subplots()
    ax_wc.imshow(wordcloud, interpolation='bilinear')
    ax_wc.axis('off')
    st.pyplot(fig_wc)

    st.header("Ranking visual de frases más interesantes")
    for i, frase in enumerate(frases_ordenadas[:5]):
        st.markdown(f"<div style='background:#A5D6A7;padding:1em;margin-bottom:0.5em;border-radius:10px;'><b>{i+1}.</b> {frase} <span style='color:gray'>({len(frase.split())} palabras)</span></div>", unsafe_allow_html=True)

    # --- Grafo interactivo de entidades ---

    # --- Infografía automática ---
    st.header("Infografía automática del documento")
    import matplotlib.pyplot as plt
    fig_info, ax_info = plt.subplots(figsize=(8, 4))
    ax_info.barh(palabras[:5], palabras_freq[:5], color='#4CAF50')
    ax_info.set_xlabel('Frecuencia')
    ax_info.set_title('Top 5 palabras clave')
    for i, v in enumerate(palabras_freq[:5]):
        ax_info.text(v + 0.1, i, str(v), color='black', va='center')
    st.pyplot(fig_info)

    # --- Guardar análisis en historial ---

    st.session_state["historial"].append({
        "nombre": archivo.name,
        "fecha": datetime.datetime.now().strftime("%d/%m/%Y %H:%M"),
        "texto": texto
    })

    # --- Exportar resultados ---
    st.header("Exportar resultados")
    resumen = '. '.join(frases_ordenadas[:5]) + '.'
    # Exportar resumen y palabras clave en PDF
    if st.button("Descargar resumen y palabras clave en PDF"):
        pdf = FPDF()
        pdf.add_page()
        pdf.set_font("Arial", size=12)
        pdf.cell(200, 10, txt="Resumen", ln=True, align='C')
        pdf.multi_cell(0, 10, resumen)
        pdf.cell(200, 10, txt="Palabras clave", ln=True, align='C')
        pdf.multi_cell(0, 10, ", ".join(palabras))
        pdf_bytes = BytesIO()
        pdf.output(pdf_bytes)
        pdf_bytes.seek(0)
        st.download_button("Descargar PDF", data=pdf_bytes, file_name="resumen_palabras.pdf", mime="application/pdf")
    # Exportar resumen y palabras clave en TXT
    if st.button("Descargar resumen y palabras clave en TXT"):
        txt = f"Resumen:\n{resumen}\n\nPalabras clave:\n{', '.join(palabras)}"
        st.download_button("Descargar TXT", data=txt, file_name="resumen_palabras.txt", mime="text/plain")
    # --- Generador automático de presentaciones ---
    st.header("Generador automático de presentación (PowerPoint)")
    from pptx import Presentation
    from pptx.util import Inches, Pt
    if st.button("Descargar presentación PPTX"):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Resumen del documento"
        slide.placeholders[1].text = resumen
        slide2 = prs.slides.add_slide(prs.slide_layouts[1])
        slide2.shapes.title.text = "Palabras clave"
        tf = slide2.shapes.placeholders[1].text_frame
        for palabra in palabras[:10]:
            tf.add_paragraph(palabra)
        # Guardar presentación
        pptx_bytes = BytesIO()
        prs.save(pptx_bytes)
        pptx_bytes.seek(0)
        st.download_button("Descargar PPTX", data=pptx_bytes, file_name="presentacion_qintella.pptx", mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")
    else:
        st.warning("No se pudo extraer texto del archivo.")

else:
    st.info("Por favor, sube un archivo para comenzar el análisis.")

# --- Configuración avanzada ---
st.sidebar.header("Configuración avanzada")
tema = st.sidebar.selectbox("Color del tema", ["Verde", "Azul", "Rojo", "Gris"])
tamano_fuente = st.sidebar.slider("Tamaño de fuente", min_value=10, max_value=30, value=14)

# Aplicar tema dinámico
if tema == "Verde":
    color_fondo = "#e6ffe6"
elif tema == "Azul":
    color_fondo = "#e6f0ff"
elif tema == "Rojo":
    color_fondo = "#ffe6e6"
else:
    color_fondo = "#f0f0f0"

st.markdown(f"""
    <style>
    .main {{background-color: {color_fondo};}}
    .stApp {{background-color: {color_fondo};}}
    .css-1d391kg {{background-color: {color_fondo};}}
    div, p, span, h1, h2, h3, h4, h5, h6 {{font-size: {tamano_fuente}px !important;}}
    </style>
""", unsafe_allow_html=True)

# Mostrar historial de análisis en la barra lateral
st.sidebar.header("Historial de análisis")
if st.session_state["historial"]:
    for i, item in enumerate(reversed(st.session_state["historial"])):
        if st.sidebar.button(f"Ver análisis: {item['nombre']} ({item['fecha']})", key=f"hist_{i}"):
            st.write(f"## Análisis previo de: {item['nombre']} ({item['fecha']})")
            st.write(item["texto"])

st.sidebar.header("Ayuda y documentación")
st.sidebar.markdown("""
**¿Cómo usar Q-Intella?**
1. Sube un archivo Word, PDF, TXT o CSV.
2. Visualiza palabras clave, resumen, gráficas y análisis de sentimiento.
3. Exporta resultados en PDF, TXT o imágenes.
4. Personaliza la interfaz en Configuración avanzada.

**Enlaces útiles:**
- [Documentación oficial de Streamlit](https://docs.streamlit.io/)
- [Soporte Python](https://www.python.org/doc/)
- [Contacto Q-Intella](mailto:soporte@q-intella.com)
""")



# --- Análisis de emociones y reconocimiento de imágenes SOLO si hay texto válido ---
if archivo and 'texto' in locals() and texto.strip():
    # --- Análisis de emociones y tono del texto ---
    st.header("Análisis de emociones y tono del texto")
    import text2emotion as t2e
    emociones = t2e.get_emotion(texto)
    emojis = {
        'Happy': '😊',
        'Angry': '😠',
        'Surprise': '😮',
        'Sad': '😢',
        'Fear': '😱'
    }
    st.write({emojis[k]: v for k, v in emociones.items()})
    st.bar_chart(emociones)

    # --- Reconocimiento de imágenes y texto en imágenes ---
    st.header("Reconocimiento de imágenes en el documento")
    import easyocr
    from PIL import Image
    import tempfile
    imagenes_encontradas = []
    # Extraer imágenes de Word
    if archivo.name.endswith(".docx"):
        docx_file = BytesIO(archivo.getvalue())
        doc = docx.Document(docx_file)
        for rel in doc.part.rels.values():
            if "image" in rel.target_ref:
                img_bytes = rel.target_part.blob
                img = Image.open(BytesIO(img_bytes))
                imagenes_encontradas.append(img)
    # Extraer imágenes de PDF
    elif archivo.name.endswith(".pdf"):
        import fitz  # PyMuPDF
        pdf_file = BytesIO(archivo.getvalue())
        pdf_doc = fitz.open(stream=pdf_file.read(), filetype="pdf")
        for page in pdf_doc:
            for img_index, img in enumerate(page.get_images(full=True)):
                xref = img[0]
                base_image = pdf_doc.extract_image(xref)
                img_bytes = base_image["image"]
                img = Image.open(BytesIO(img_bytes))
                imagenes_encontradas.append(img)
    # Mostrar y analizar imágenes
    if imagenes_encontradas:
        st.subheader(f"Imágenes encontradas: {len(imagenes_encontradas)}")
        reader = easyocr.Reader(['es', 'en'])
        for idx, img in enumerate(imagenes_encontradas):
            st.image(img, caption=f"Imagen {idx+1}", use_column_width=True)
            st.write("Texto detectado en la imagen:")
            resultado = reader.readtext(np.array(img))
            texto_detectado = " ".join([r[1] for r in resultado])
            st.write(texto_detectado if texto_detectado else "No se detectó texto.")
    else:
        st.info("No se encontraron imágenes en el documento.")

elif archivo and 'texto' in locals() and not texto.strip():
    st.warning("No se pudo extraer texto del archivo.")
else:
    st.info("Por favor, sube un archivo para comenzar el análisis.")
    